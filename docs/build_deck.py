"""
Generates the Playwright Enterprise E2E Test Framework architecture deck as a
native .pptx (imports cleanly into Google Slides).

Run:  python docs/build_deck.py
Out:  docs/Playwright-E2E-Framework-Architecture.pptx

Theme: Navy + Teal. 12 slides, 16:9, with speaker notes.
Re-runnable — edit the content blocks below and regenerate.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image as PILImage
import os

# ---------------------------------------------------------------- palette
NAVY   = RGBColor(0x0F, 0x2D, 0x4A)
NAVY2  = RGBColor(0x1B, 0x41, 0x66)
TEAL   = RGBColor(0x2B, 0xB2, 0xA6)
INK    = RGBColor(0x1F, 0x29, 0x37)
GRAY   = RGBColor(0x5B, 0x64, 0x72)
LGRAY  = RGBColor(0x94, 0x9C, 0xA8)
LIGHT  = RGBColor(0xF3, 0xF6, 0xF9)
CARD   = RGBColor(0xF7, 0xF9, 0xFB)
BORDER = RGBColor(0xD8, 0xDE, 0xE6)
GREEN  = RGBColor(0x16, 0xA3, 0x4A)
AMBER  = RGBColor(0xF1, 0x9E, 0x0B)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Segoe UI"
SW, SH = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.6)
CONTENT_W = Inches(13.333 - 1.2)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

_page = 0

# ---------------------------------------------------------------- helpers
def _set(run, size, color, bold=False, italic=False, font=FONT):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font

def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tf

def para(tf, text, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT,
         space_after=6, space_before=0, level=0, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    p.level = level
    r = p.add_run()
    r.text = text
    _set(r, size, color, bold, italic)
    return p

def rect(slide, l, t, w, h, fill, line=None, line_w=0.75, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, l, t, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp

def shape_text(sp, lines, default_size=13):
    """lines: list of (text, size, color, bold, align)."""
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    for i, (text, size, color, bold, align) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(2)
        r = p.add_run()
        r.text = text
        _set(r, size, color, bold)

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def base_slide(title, kicker=None):
    global _page
    _page += 1
    slide = prs.slides.add_slide(BLANK)
    # title bar
    rect(slide, 0, 0, SW, Inches(0.95), NAVY)
    rect(slide, 0, Inches(0.95), SW, Inches(0.06), TEAL)
    tf = textbox(slide, MARGIN, Inches(0.16), Inches(11.8), Inches(0.7), MSO_ANCHOR.MIDDLE)
    if kicker:
        para(tf, kicker, 11, TEAL, bold=True, space_after=1, first=True)
        para(tf, title, 24, WHITE, bold=True, space_after=0)
    else:
        para(tf, title, 26, WHITE, bold=True, space_after=0, first=True)
    # footer
    ftf = textbox(slide, MARGIN, Inches(7.08), Inches(9), Inches(0.3))
    para(ftf, "Playwright Enterprise E2E Test Framework", 9, LGRAY, first=True)
    ptf = textbox(slide, Inches(11.6), Inches(7.08), Inches(1.13), Inches(0.3))
    para(ptf, str(_page), 9, LGRAY, align=PP_ALIGN.RIGHT, first=True)
    return slide

def bullets(slide, items, l=MARGIN, t=Inches(1.35), w=CONTENT_W, h=Inches(5.4), size=16):
    """items: list of (text, level). level 0 = •, level 1 = –."""
    tf = textbox(slide, l, t, w, h)
    for i, (text, level) in enumerate(items):
        glyph = "•  " if level == 0 else "–  "
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(7 if level == 0 else 4)
        p.level = level
        r = p.add_run()
        r.text = glyph + text
        _set(r, size if level == 0 else size - 2,
             INK if level == 0 else GRAY, bold=False)
    return tf

def card(slide, l, t, w, h, title, lines, accent=TEAL):
    box = rect(slide, l, t, w, h, CARD, line=BORDER, line_w=0.75,
               shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(slide, l, t, Inches(0.09), h, accent)  # left accent bar
    tf = textbox(slide, l + Inches(0.28), t + Inches(0.16), w - Inches(0.42),
                 h - Inches(0.3))
    para(tf, title, 15, NAVY, bold=True, space_after=5, first=True)
    for ln in lines:
        para(tf, ln, 11.5, GRAY, space_after=3)
    return box

def badge(slide, l, t, text, color, w=Inches(1.2), h=Inches(0.34)):
    b = rect(slide, l, t, w, h, color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    shape_text(b, [(text, 10.5, WHITE, True, PP_ALIGN.CENTER)])
    return b

def arrow(slide, l, t, w, h=Inches(0.4)):
    rect(slide, l, t, w, h, TEAL, shape=MSO_SHAPE.RIGHT_ARROW)

def comp_table(slide, rows, l, t, w, col_widths, header_h=Inches(0.48),
               row_h=Inches(0.52), accent_col=2):
    n_rows, n_cols = len(rows), len(rows[0])
    total_h = header_h + row_h * (n_rows - 1)
    gt = slide.shapes.add_table(n_rows, n_cols, l, t, w, total_h).table
    gt.first_row = False
    gt.horz_banding = False
    for ci, cw in enumerate(col_widths):
        gt.columns[ci].width = cw
    for ri, row in enumerate(rows):
        gt.rows[ri].height = header_h if ri == 0 else row_h
        for cix, val in enumerate(row):
            cell = gt.cell(ri, cix)
            cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            r = p.add_run(); r.text = val
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = TEAL if cix == accent_col else NAVY
                _set(r, 12, WHITE, bold=True)
            else:
                cell.fill.fore_color.rgb = WHITE if ri % 2 == 1 else LIGHT
                if cix == 0:
                    _set(r, 11.5, NAVY, bold=True)
                else:
                    _set(r, 11.5, INK if cix == accent_col else GRAY,
                         bold=(cix == accent_col))
    return gt

def phase_chip(slide, l, t, w, num, title, desc):
    box = rect(slide, l, t, w, Inches(0.92), CARD, line=BORDER, line_w=0.75,
               shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(slide, l, t, Inches(0.09), Inches(0.92), TEAL)
    tf = textbox(slide, l + Inches(0.24), t + Inches(0.1), w - Inches(0.36), Inches(0.75))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = num + "  "; _set(r, 14, TEAL, bold=True)
    r2 = p.add_run(); r2.text = title; _set(r2, 12.5, NAVY, bold=True)
    para(tf, desc, 10.5, GRAY, space_before=2)

def add_image(slide, img_path, left, top, max_w, max_h, border=True):
    """Places an image scaled to fit within (max_w x max_h), preserving aspect,
    top-left anchored at (left, top). Returns (width, height) in EMU."""
    iw, ih = PILImage.open(img_path).size
    ratio = iw / ih
    width = int(max_w)
    height = int(width / ratio)
    if height > int(max_h):
        height = int(max_h)
        width = int(height * ratio)
    pic = slide.shapes.add_picture(img_path, left, top, width=width, height=height)
    if border:
        pic.line.color.rgb = BORDER
        pic.line.width = Pt(0.75)
    return width, height

IMG_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "img")

# ================================================================ SLIDE 1 — Title
def s_title():
    global _page
    _page += 1
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, SW, SH, NAVY)
    rect(slide, 0, Inches(3.9), SW, Inches(0.08), TEAL)
    tf = textbox(slide, Inches(0.9), Inches(2.35), Inches(11.5), Inches(1.6))
    para(tf, "PLAYWRIGHT ENTERPRISE E2E TEST FRAMEWORK", 13, TEAL, bold=True,
         space_after=8, first=True)
    para(tf, "Architecture & CI/CD Overview", 40, WHITE, bold=True, space_after=0)
    tf2 = textbox(slide, Inches(0.9), Inches(4.2), Inches(11.5), Inches(1.6))
    para(tf2, "Page Object Model  ·  Data-Driven  ·  Multi-Auth API  ·  Rich Reporting  ·  GitHub Actions CI/CD",
         15, RGBColor(0xC9, 0xD6, 0xE2), space_after=18, first=True)
    para(tf2, "Prepared for: Technical Architecture Review", 13, WHITE, bold=True, space_after=2)
    para(tf2, "Presenter: Gukan  ·  QA Engineering", 12, LGRAY)
    notes(slide,
          "One-line intro: this is our in-house Playwright + TypeScript E2E framework. "
          "Goal of this review: make the case for the framework vs. a bare test set, walk the "
          "architecture and the reporting/CI-CD pipeline, and show how we fold the existing "
          "web-pet-golang Playwright e2e specs in without duplicating infrastructure.")

# ================================================================ SLIDE 2 — Migration plan / why this framework
def s_migration():
    slide = base_slide("Migration Plan — Why This Framework", kicker="THE CASE")
    tf = textbox(slide, MARGIN, Inches(1.15), CONTENT_W, Inches(0.45))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "The web-pet e2e suite is a set of tests — not a framework.  "
    _set(r, 15, INK, bold=True)
    r2 = p.add_run(); r2.text = "This project is the platform those tests should run on."
    _set(r2, 15, TEAL, bold=True)

    stf = textbox(slide, MARGIN, Inches(1.6), CONTENT_W, Inches(0.32))
    para(stf, "Scope: this targets the Playwright e2e layer only — 56 specs. The web-pet Go API (451) "
              "and Vitest (230) test layers are separate and unaffected.", 11.5, GRAY, first=True)

    rows = [
        ["Capability", "web-pet e2e  (tests only)", "This framework"],
        ["Reusable POM / components", "No page-object layer", "BasePage + Nav / Modal / Form"],
        ["External test data", "By-name id lookups in specs", "JSON / CSV data-driven layer"],
        ["API authentication", "Session login (storageState)", "OAuth2 / Basic / API-Key + retry"],
        ["Reporting", "Playwright HTML only", "HTML · Allure · Email · Slack · ELK*"],
        ["CI/CD & notifications", "Manual run; no notify pipeline", "Push / cron / dispatch + Email / Slack"],
    ]
    comp_table(slide, rows, MARGIN, Inches(1.98), CONTENT_W,
               [Inches(3.0), Inches(4.5), Inches(4.63)],
               header_h=Inches(0.46), row_h=Inches(0.48))

    ntf = textbox(slide, MARGIN, Inches(4.88), CONTENT_W, Inches(0.3))
    para(ntf, "web-pet's strengths — storageState auth, resolve-ids-by-name, and the equiv/ legacy-parity "
              "suite — are preserved as framework conventions.   * ELK in progress.",
         10.5, LGRAY, italic=True, first=True)

    ltf = textbox(slide, MARGIN, Inches(5.22), CONTENT_W, Inches(0.3))
    para(ltf, "MIGRATION PLAN", 11, TEAL, bold=True, first=True)
    cw = Inches(2.95)
    gap = Inches(0.11)
    phases = [
        ("1", "Adopt the harness", "Point specs at framework fixtures, config & auth; retire parallel glue."),
        ("2", "Refactor to POM", "Pickers → components; shared helpers; drop ad-hoc selectors."),
        ("3", "Wire reporting & CI", "Route through Allure / Email / Slack; run under CI triggers."),
        ("4", "Preserve parity", "Keep equiv/ + by-name ids; then broaden coverage."),
    ]
    x = MARGIN
    for num, title, desc in phases:
        phase_chip(slide, x, Inches(5.52), cw, num, title, desc)
        x = x + cw + gap
    notes(slide, "Open with the core argument: the web-pet Playwright layer is 56 spec files plus minimal "
                 "glue — no reusable engineering layer (no POM base, no external data layer, no multi-auth "
                 "API, no rich reporting, no CI notification pipeline). Our project IS that layer. The move "
                 "is not 'rebuild tests' — it's migrate the existing specs onto this platform and delete the "
                 "duplicate glue. Scope only the e2e layer; the Go (451) and Vitest (230) layers stay put. "
                 "Four phases; and be fair — web-pet's storageState auth and equiv/ parity are good, we keep them.")

# ================================================================ SLIDE 3 — Agenda
def s_agenda():
    slide = base_slide("Agenda")
    items = [
        ("Migration plan & why this framework — tests vs. platform", 0),
        ("Framework at a glance — purpose, stack, design principles", 0),
        ("Architecture — the layered design", 0),
        ("Project structure — how the code is organised", 0),
        ("Core capabilities — POM, data-driven, fixtures, API auth, observability", 0),
        ("AI-augmented authoring — Claude skills & Playwright agents", 0),
        ("Reporting — Allure size-fix, rich Email & Slack, ELK (in progress)", 0),
        ("CI/CD — triggers (dev push, cron, manual), pipeline, notifications, environments", 0),
        ("Consolidating the web-pet-golang e2e suite — one harness, no duplication", 0),
        ("Roadmap & status", 0),
    ]
    bullets(slide, items, t=Inches(1.5), size=18)
    notes(slide, "Keep this quick — it's a signpost. Emphasise two things the architect cares "
                 "about most: the CI/CD trigger model and the no-duplication consolidation of the "
                 "existing 56 web-pet specs.")

# ================================================================ SLIDE 3 — At a glance
def s_glance():
    slide = base_slide("At a Glance")
    top = Inches(1.35)
    cw, gap = Inches(3.83), Inches(0.22)
    card(slide, MARGIN, top, cw, Inches(2.2), "What it is", [
        "Enterprise Playwright + TypeScript E2E framework for the PET Tiger app.",
        "Page Object Model + component architecture.",
        "Scales with coverage — no upfront over-engineering.",
    ])
    card(slide, MARGIN + cw + gap, top, cw, Inches(2.2), "Tech stack", [
        "Playwright Test 1.58  ·  TypeScript 5",
        "Allure reporting + trend history",
        "GitHub Actions CI  ·  Node 22 / Java 21",
        "Nodemailer (email)  ·  Slack webhook",
    ])
    card(slide, MARGIN + 2 * (cw + gap), top, cw, Inches(2.2), "Consolidation scope", [
        "web-pet-golang has 3 test layers: Go API (451), Vitest (230), Playwright e2e (56).",
        "This platform consolidates the Playwright e2e layer only.",
        "Go & Vitest layers are unaffected.",
    ], accent=NAVY2)
    # value band
    band = rect(slide, MARGIN, Inches(3.95), CONTENT_W, Inches(1.0), NAVY,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    shape_text(band, [
        ("One harness: authored tests + data + auth + reporting + CI, in a single maintained codebase.",
         16, WHITE, True, PP_ALIGN.CENTER),
    ])
    tf = textbox(slide, MARGIN, Inches(5.2), CONTENT_W, Inches(1.6))
    para(tf, "Design principle", 13, TEAL, bold=True, space_after=4, first=True)
    para(tf, "Thin base classes over native Playwright APIs — no redundant wrappers; folders earn "
             "their place by real usage. Recent cleanup removed ~40 unused wrapper methods and pruned "
             "unused type surface to keep the framework lean and navigable.", 13, GRAY)
    notes(slide, "Frame it as intentional and disciplined: we deliberately keep base classes thin "
                 "(call Playwright's Locator/expect directly) and only add structure when a real "
                 "second consumer appears. That keeps onboarding fast.")

# ================================================================ SLIDE 4 — Architecture diagram
def s_architecture():
    slide = base_slide("Architecture", kicker="LAYERED DESIGN")
    layers = [
        ("L1", "CI / CD & Triggers", "GitHub Actions — push, cron, manual, external dispatch"),
        ("L2", "Configuration", "playwright.config.ts  ·  envLoader  ·  env.local / dev / qa"),
        ("L3", "Global Lifecycle", "global-setup (auth dir, results)  ·  global-teardown"),
        ("L4", "Fixtures & Hooks", "base.fixture / api.fixture  ·  beforeEach / afterEach  ·  lifecycle manager"),
        ("L5", "Page Objects & Components", "BasePage  ·  BaseComponent  ·  Navigation / Modal / Form"),
        ("L6", "Data Layer", "DataProvider  ·  JSON / CSV readers (read directly, no conversion)"),
        ("L7", "Authentication & API", "AuthorizationManager  ·  RequestBuilder  ·  AuthContextFactory"),
        ("L8", "Reporting & Utilities", "HTML · Allure · Email · Slack · ELK  ·  Logger · network helpers"),
    ]
    top = Inches(1.28)
    row_h = Inches(0.63)
    gap = Inches(0.055)
    for i, (tag, name, desc) in enumerate(layers):
        y = top + Emu(int((row_h + gap) * i))
        bar = rect(slide, MARGIN, y, CONTENT_W, row_h, CARD, line=BORDER, line_w=0.5,
                   shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        chip = rect(slide, MARGIN, y, Inches(0.85), row_h, NAVY,
                    shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        shape_text(chip, [(tag, 14, WHITE, True, PP_ALIGN.CENTER)])
        tf = textbox(slide, MARGIN + Inches(1.05), y, Inches(3.6), row_h, MSO_ANCHOR.MIDDLE)
        para(tf, name, 14, NAVY, bold=True, space_after=0, first=True)
        tf2 = textbox(slide, MARGIN + Inches(4.8), y, Inches(7.2), row_h, MSO_ANCHOR.MIDDLE)
        para(tf2, desc, 11.5, GRAY, space_after=0, first=True)
    notes(slide, "Walk top to bottom: a run is triggered (L1), config resolves the environment (L2), "
                 "global setup logs in once and saves storage state (L3), fixtures inject page objects "
                 "and run lifecycle hooks (L4), tests drive pages/components (L5) with data from JSON/CSV "
                 "(L6), API calls go through the auth layer (L7), and everything fans out to reporting (L8).")

# ================================================================ SLIDE 5 — Project structure
def s_structure():
    slide = base_slide("Project Structure", kicker="src/")
    cols = [
        ("core/", "Framework primitives — constants, exceptions, @HandleError decorator"),
        ("config/", "Env loading, environment definitions, data-source resolution"),
        ("auth/", "OAuth2 / Basic / API-Key request auth + 401/403 retry"),
        ("components/", "BaseComponent + Navigation / Modal / Form (root-scoped)"),
        ("context/", "executionContext (runId/branch/commit) · testMetrics · testRunContext"),
        ("data/", "runnerManager.json / .csv · module data · optional runnerList filter"),
        ("fixtures/", "base.fixture (UI) · api.fixture · global setup / teardown"),
        ("pages/", "BasePage + LoginPage / LeftNavigationPage"),
        ("reporting/", "email · slack · dashboard (ELK) · shared runSummary"),
        ("utils/", "DataProvider · logger · allureHelper · network · assertions · dataReaders/"),
    ]
    top = Inches(1.35)
    cw = Inches(6.0)
    ch = Inches(0.52)
    gap = Inches(0.08)
    for i, (name, desc) in enumerate(cols):
        col = i % 2
        row = i // 2
        l = MARGIN + col * (cw + Inches(0.15))
        t = top + Emu(int((ch + gap) * row))
        b = rect(slide, l, t, cw, ch, CARD, line=BORDER, line_w=0.5,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tf = b.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.14); tf.margin_right = Inches(0.1)
        tf.margin_top = 0; tf.margin_bottom = 0
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = name + "   "
        _set(r, 13, TEAL, bold=True)
        r2 = p.add_run(); r2.text = desc
        _set(r2, 11, GRAY)
    tf = textbox(slide, MARGIN, Inches(6.55), CONTENT_W, Inches(0.5))
    para(tf, "tests/  — auth.setup.ts (storageState login)  ·  login/  ·  spec files (**/*.spec.ts)",
         12, INK, bold=True, first=True)
    notes(slide, "Point out the recent tidy-up: single-file concerns collapsed into core/, and the "
                 "structure is flat where it should be. Each folder maps to a capability we'll cover next.")

# ================================================================ SLIDE 6 — Capabilities I
def s_cap1():
    slide = base_slide("Core Capabilities  (1 / 2)", kicker="THE FRAMEWORK")
    top = Inches(1.35)
    cw = Inches(5.95)
    card(slide, MARGIN, top, cw, Inches(2.55), "Page Object Model & Components", [
        "BasePage — thin: navigate(), waitForCondition(), screenshot convention.",
        "Actions/assertions use native Locator + expect() directly (no wrappers).",
        "BaseComponent — root-scoped locator finders; Navigation / Modal / Form.",
        "Prevents selector collisions; pages stay small and readable.",
    ])
    card(slide, MARGIN + cw + Inches(0.2), top, cw, Inches(2.55), "Data-Driven Testing", [
        "Reads JSON or CSV directly — no conversion / preprocessing step.",
        "DataProvider singleton + JSON/CSV readers behind one interface.",
        "testCaseId / testCaseName fixtures auto-load & validate a row.",
        "Disabled rows (enabled:false) auto-skip.",
    ], accent=NAVY2)
    band = rect(slide, MARGIN, Inches(4.2), CONTENT_W, Inches(1.7), LIGHT,
                line=BORDER, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = textbox(slide, MARGIN + Inches(0.3), Inches(4.4), CONTENT_W - Inches(0.6), Inches(1.4))
    para(tf, "Why it matters", 13, TEAL, bold=True, space_after=4, first=True)
    para(tf, "New pages plug into a proven base; new scenarios are added as data rows, not new code. "
             "This is how the suite scales without the page layer ballooning.", 13, GRAY)
    notes(slide, "POM keeps locators in one place per page; components keep reusable widgets DRY. "
                 "Data-driven means QA can add cases by editing a JSON/CSV row — no TS changes — which "
                 "is the key to scaling coverage.")

# ================================================================ SLIDE 7 — Capabilities II
def s_cap2():
    slide = base_slide("Core Capabilities  (2 / 2)", kicker="THE FRAMEWORK")
    top = Inches(1.35)
    cw = Inches(3.83)
    gap = Inches(0.22)
    card(slide, MARGIN, top, cw, Inches(2.7), "Fixtures & Lifecycle", [
        "base.fixture (UI) & api.fixture (API-only).",
        "beforeEach/afterEach → lifecycle manager.",
        "TestMetrics per test (status, duration, retry).",
        "ExecutionContext: runId, branch, commit, trigger.",
    ])
    card(slide, MARGIN + cw + gap, top, cw, Inches(2.7), "API Testing & Auth", [
        "Strategies: OAuth2, Basic, API-Key, none.",
        "OAuth2 token cache + auto-refresh.",
        "executeWithAuthRetry — retry once on 401/403.",
        "authContextFactory builds the request context.",
    ], accent=NAVY2)
    card(slide, MARGIN + 2 * (cw + gap), top, cw, Inches(2.7), "Logging & Observability", [
        "Structured Logger — console + daily file.",
        "Per-test metrics + run-level context.",
        "Screenshot / video / trace on failure.",
        "Feeds the reporting layer.",
    ])
    tf = textbox(slide, MARGIN, Inches(4.35), CONTENT_W, Inches(1.6))
    para(tf, "Browser auth = saved storageState", 13, TEAL, bold=True, space_after=4, first=True)
    para(tf, "A one-time auth-setup project logs in and persists the session to .auth/user.json; every "
             "browser test starts authenticated. API auth is separate and configurable per environment.",
         13, GRAY)
    notes(slide, "Two independent auth paths: (1) browser login once via storageState for UI tests, "
                 "(2) API request auth (OAuth2/Basic/API-Key) for direct API calls and validation. "
                 "ExecutionContext gives every report the run's branch/commit/trigger.")

# ================================================================ SLIDE — AI-augmented authoring
def s_ai():
    slide = base_slide("AI-Augmented Test Engineering", kicker="CLAUDE CODE")
    itf = textbox(slide, MARGIN, Inches(1.12), CONTENT_W, Inches(0.4))
    para(itf, "Tests are authored, planned and repaired with Claude Code. Skills encode THIS framework's "
              "conventions; agents drive a real browser via MCP — so output follows the POM, fixtures and "
              "data-driven layer and runs without rework.", 12.5, INK, bold=True, first=True)

    def kv(tf, name, desc, first=False):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.space_after = Pt(6)
        r = p.add_run(); r.text = name + "  "; _set(r, 11.5, TEAL, bold=True)
        r2 = p.add_run(); r2.text = "— " + desc; _set(r2, 11, GRAY)

    top = Inches(1.72)
    cw = Inches(5.95)
    ch = Inches(3.35)
    # Left — Claude Skills
    rect(slide, MARGIN, top, cw, ch, CARD, line=BORDER, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(slide, MARGIN, top, Inches(0.09), ch, TEAL)
    ltf = textbox(slide, MARGIN + Inches(0.28), top + Inches(0.18), cw - Inches(0.5), ch - Inches(0.35))
    para(ltf, "Claude Skills  (repo-scoped)", 14, NAVY, bold=True, space_after=8, first=True)
    kv(ltf, "ui-script-generator", "generate/update a UI spec from a chat scenario; encodes POM, fixtures, data-driven & locator rules.")
    kv(ltf, "data-driven-testing", "runnerManager rows + testCaseId / testCaseName fixtures; JSON ↔ CSV, no conversion.")
    kv(ltf, "tdd", "red → green: write the failing spec from acceptance criteria first, then build page objects.")
    kv(ltf, "jira-to-script", "Jira ticket → plan → generate → run → heal (via Jira MCP + the agents).")

    # Right — Playwright Agents
    rx = MARGIN + cw + Inches(0.2)
    rect(slide, rx, top, cw, ch, CARD, line=BORDER, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(slide, rx, top, Inches(0.09), ch, NAVY2)
    rtf = textbox(slide, rx + Inches(0.28), top + Inches(0.18), cw - Inches(0.5), ch - Inches(0.35))
    para(rtf, "Playwright Agents  (subagents)", 14, NAVY, bold=True, space_after=8, first=True)
    kv(rtf, "planner", "explores the live app in a real browser and saves a comprehensive test plan.")
    kv(rtf, "generator", "executes each planned step live, then writes the Playwright spec.")
    kv(rtf, "healer", "runs failing tests; debugs via snapshot / console / network; fixes selectors & timing.")
    para(rtf, "Each runs as an isolated subagent with only its Playwright MCP tools — planner & generator "
              "author, healer repairs.", 10.5, LGRAY, italic=True, space_before=6)

    # Bottom band — MCP
    by = Inches(5.28)
    rect(slide, MARGIN, by, CONTENT_W, Inches(1.15), NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    btf = textbox(slide, MARGIN + Inches(0.32), by + Inches(0.15), CONTENT_W - Inches(0.64), Inches(0.95),
                  MSO_ANCHOR.MIDDLE)
    para(btf, "Powered by MCP (.mcp.json)", 12.5, TEAL, bold=True, space_after=3, first=True)
    para(btf, "playwright-test server → real-browser tools (snapshot, click, generate-locator, run/debug) · "
              "jira server → pulls stories.  The skills carry the repo's conventions, so generated specs use "
              "the POM / fixtures / data-driven layer and pass without manual correction.", 11.5, WHITE)

    notes(slide, "This is the differentiator: authoring isn't manual. Four repo-scoped Claude skills encode "
                 "our conventions — ui-script-generator (scenario in chat → spec), data-driven-testing "
                 "(runnerManager rows + fixtures), tdd (red/green), and jira-to-script (ticket → passing "
                 "spec). Three Playwright subagents do the browser work via the playwright-test MCP server: "
                 "the planner explores the live app and writes a plan, the generator executes each step live "
                 "then emits the spec, and the healer runs failures and repairs selectors/timing. The key "
                 "point: because the skills embed the framework's POM/fixtures/data-driven rules, the AI "
                 "output drops straight into the suite and runs — it doesn't fight our conventions.")

# ================================================================ SLIDE — Reporting
def s_reporting():
    slide = base_slide("Reporting", kicker="OUTPUTS")
    itf = textbox(slide, MARGIN, Inches(1.12), CONTENT_W, Inches(0.34))
    para(itf, "Environment-tagged, rich notifications — real Email & Slack output shown below.",
         13, INK, bold=True, first=True)

    # Email image (left) + caption
    etf = textbox(slide, MARGIN, Inches(1.52), Inches(5.5), Inches(0.28))
    para(etf, "EMAIL NOTIFICATION", 10.5, TEAL, bold=True, first=True)
    add_image(slide, os.path.join(IMG_DIR, "email-report.png"),
              MARGIN, Inches(1.82), Inches(5.35), Inches(4.55))

    # Slack image (right) + caption
    rx = Inches(6.35)
    stf = textbox(slide, rx, Inches(1.52), Inches(6.4), Inches(0.28))
    para(stf, "SLACK NOTIFICATION", 10.5, TEAL, bold=True, first=True)
    add_image(slide, os.path.join(IMG_DIR, "slack-report.png"),
              rx, Inches(1.82), Inches(6.35), Inches(2.7))

    # Highlights card (right, under slack)
    hy = Inches(4.72)
    rect(slide, rx, hy, Inches(6.35), Inches(1.75), NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    htf = textbox(slide, rx + Inches(0.28), hy + Inches(0.16), Inches(5.8), Inches(1.5))
    para(htf, "Highlights", 13, TEAL, bold=True, space_after=5, first=True)
    para(htf, "• Allure embeds screenshots only → the size failure is gone (~3 MB vs 5.8 MB); "
              "video & trace stay in the Playwright HTML report + artifacts.", 11.5, WHITE, space_after=4)
    para(htf, "• Env badges: LOCAL / DEV / QA + CI · pass rate · failures with spec & error · run link.",
         11.5, WHITE, space_after=4)
    para(htf, "• Channels: List · HTML · JSON · Allure · Email · Slack · ELK (in progress). "
              "Email/Slack/ELK share one runSummary and are self-gating (opt-in).",
         11.5, RGBColor(0xC9, 0xD6, 0xE2))
    badge(slide, rx + Inches(4.95), hy + Inches(0.16), "ELK: WIP", AMBER, w=Inches(1.15), h=Inches(0.3))

    notes(slide, "Show, don't tell: these are our actual rich Email and Slack notifications, "
                 "environment-tagged (DEV + CI badges), with pass rate, counts, duration, metadata, and "
                 "failures showing the spec file + error line. The key reliability win is the Allure "
                 "size-fix — screenshots only, so the single-file report no longer blows the size limit; "
                 "video and full trace remain in the Playwright HTML report and CI artifacts. Email, Slack "
                 "and ELK share one runSummary module; all are opt-in. ELK is scaffolded, not yet wired live.")

# ================================================================ SLIDE 9 — CI/CD pipeline & triggers
def s_cicd_flow():
    slide = base_slide("CI/CD Pipeline & Triggers", kicker="GITHUB ACTIONS")
    itf = textbox(slide, MARGIN, Inches(1.1), CONTENT_W, Inches(0.32))
    para(itf, "End-to-end: a developer commit flows all the way to Developer & QA reviewing the results.",
         12.5, INK, bold=True, first=True)

    box_w = Inches(2.12); box_h = Inches(0.95); gap = Inches(0.28)
    arr_h = Inches(0.36)
    row1_y = Inches(1.65); row2_y = Inches(3.42)

    def colx(i):
        return MARGIN + i * (box_w + gap)

    def fbox(x, y, text, fill):
        b = rect(slide, x, y, box_w, box_h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        shape_text(b, [(text, 11, WHITE, True, PP_ALIGN.CENTER)])

    # Row 1 (left → right): steps 1-5
    row1 = [("Developer\nCommit", NAVY), ("Dev\nRepository", NAVY), ("GitHub\nActions", NAVY),
            ("Build", NAVY), ("Deploy\nto DEV", NAVY2)]
    for i, (txt, fill) in enumerate(row1):
        fbox(colx(i), row1_y, txt, fill)
        if i < 4:
            rect(slide, colx(i) + box_w, row1_y + (box_h - arr_h) // 2, gap, arr_h, TEAL,
                 shape=MSO_SHAPE.RIGHT_ARROW)

    # Down arrow: Deploy to DEV → Trigger Playwright (the key handoff)
    rect(slide, colx(4) + box_w // 2 - Inches(0.22), row1_y + box_h, Inches(0.44),
         row2_y - (row1_y + box_h), TEAL, shape=MSO_SHAPE.DOWN_ARROW)

    # Row 2 (right → left, serpentine): steps 6-10
    row2 = [("Trigger Playwright\nAutomation", TEAL), ("Run Playwright\nTests", TEAL),
            ("Generate Report\n& Artifacts", TEAL), ("Return\nPASS / FAIL", NAVY2),
            ("Developer & QA\nReview Results", NAVY)]
    for k, (txt, fill) in enumerate(row2):
        col = 4 - k
        fbox(colx(col), row2_y, txt, fill)
        if col > 0:  # left-pointing arrow into the next (col-1) box
            rect(slide, colx(col - 1) + box_w, row2_y + (box_h - arr_h) // 2, gap, arr_h, TEAL,
                 shape=MSO_SHAPE.LEFT_ARROW)

    # Cron / manual trigger feeding into Trigger Playwright Automation (col 4)
    cron_y = row2_y + box_h + Inches(0.34)
    rect(slide, colx(4) + box_w // 2 - Inches(0.16), row2_y + box_h, Inches(0.32),
         Inches(0.34), AMBER, shape=MSO_SHAPE.UP_ARROW)
    cbox = rect(slide, colx(4), cron_y, box_w, Inches(0.62), AMBER,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    shape_text(cbox, [("Cron 4:00 PM IST", 10.5, WHITE, True, PP_ALIGN.CENTER),
                      ("+ manual dispatch", 9.5, WHITE, False, PP_ALIGN.CENTER)])

    # Triggers / results note band
    band_y = Inches(5.5)
    rect(slide, MARGIN, band_y, CONTENT_W, Inches(0.98), LIGHT, line=TEAL, line_w=1.25,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = textbox(slide, MARGIN + Inches(0.3), band_y + Inches(0.13), CONTENT_W - Inches(0.6),
                 Inches(0.78), MSO_ANCHOR.MIDDLE)
    para(tf, "Two entry points into the automation:  (1) dev push → Build → Deploy to DEV auto-triggers it;  "
             "(2) Cron 4:00 PM IST daily / manual dispatch.", 12.5, INK, bold=True, first=True)
    para(tf, "Results — PASS / FAIL plus Allure, Email & Slack — flow back to Developer & QA.",
         11.5, GRAY, space_before=2)

    ptf = textbox(slide, MARGIN, Inches(6.6), CONTENT_W, Inches(0.3))
    para(ptf, "CI policy: workers=1 (shared storageState) · retries=2 · test.only blocked · concurrency-cancel per ref.",
         10.5, LGRAY, first=True)

    notes(slide, "Walk the flow: a developer commit lands in the dev repo, GitHub Actions builds and deploys "
                 "to DEV, and the DEV deploy automatically triggers the Playwright automation — it runs the "
                 "tests, generates the report & artifacts, returns a PASS/FAIL status, and developer & QA "
                 "review the results. The amber path shows the second entry point: a 4:00 PM IST cron (and "
                 "manual dispatch) that triggers the automation directly, independent of a deploy. Results go "
                 "back via Allure, Email and Slack.")

# ================================================================ SLIDE 10 — CI/CD features & environments
def s_cicd_features():
    slide = base_slide("CI/CD Features, Notifications & Environments", kicker="INFRASTRUCTURE")
    top = Inches(1.35)
    cw = Inches(3.83)
    gap = Inches(0.22)
    card(slide, MARGIN, top, cw, Inches(2.8), "Artifacts & History", [
        "Playwright HTML report (with traces).",
        "Allure report — screenshots only.",
        "Trend history cached across runs.",
        "Retained as CI artifacts per run.",
    ])
    card(slide, MARGIN + cw + gap, top, cw, Inches(2.8), "Notifications", [
        "Email — rich HTML + lean Allure attach.",
        "Slack — Block Kit summary + color bar.",
        "ELK push — in progress.",
        "All self-gating (SEND_* env flags).",
    ], accent=NAVY2)
    card(slide, MARGIN + 2 * (cw + gap), top, cw, Inches(2.8), "Environments & Secrets", [
        "env.local / env.dev / env.qa  ·  TEST_ENV.",
        "CI secrets take precedence over files.",
        "Env badges surface where a run executed.",
        "SMTP / Slack / creds via CI secrets only.",
    ])
    tf = textbox(slide, MARGIN, Inches(4.45), CONTENT_W, Inches(1.5))
    para(tf, "Secrets never live in the repo", 13, TEAL, bold=True, space_after=4, first=True)
    para(tf, "Committed env files hold only non-secret defaults; real credentials (app login, SMTP, Slack "
             "webhook) come from CI secret stores or the local shell, with OS/CI values always winning.",
         13, GRAY)
    notes(slide, "Reassure on security: no secrets in git; CI secrets override committed defaults. "
                 "Notifications are opt-in per channel via SEND_EMAIL / SEND_SLACK / SEND_RESULT_ELK, "
                 "so local runs stay silent unless you ask.")

# ================================================================ SLIDE 11 — web-pet consolidation
def s_consolidation():
    slide = base_slide("Consolidating the web-pet e2e Suite", kicker="ONE HARNESS · NO DUPLICATION")
    top = Inches(1.28)
    cw = Inches(5.95)
    card(slide, MARGIN, top, cw, Inches(2.75), "web-pet e2e today (56 specs)", [
        "Chromium-only; import { test } from ./fixtures (storageState auth).",
        "global-setup logs in once; restricted-user gating for scope tests.",
        "parent-picker-helpers.ts for non-native pickers; resolve ids by name.",
        "equiv/ — 7 parity specs (scenario.yaml, RUN_TOKEN, env-gated).",
        "No page-object layer; kebab-case + JSDoc headers; onBlur gotchas.",
    ], accent=NAVY2)
    card(slide, MARGIN + cw + Inches(0.2), top, cw, Inches(2.75), "Duplication risk", [
        "Two parallel fixture layers.",
        "Two storageState auth setups.",
        "Two Playwright configs & reporting paths.",
        "Divergent conventions over time.",
    ], accent=AMBER)
    band = rect(slide, MARGIN, Inches(4.28), CONTENT_W, Inches(2.12), NAVY,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = textbox(slide, MARGIN + Inches(0.32), Inches(4.46), CONTENT_W - Inches(0.64), Inches(1.8))
    para(tf, "Consolidation strategy — specs move into the framework, reuse one harness", 14, TEAL,
         bold=True, space_after=6, first=True)
    para(tf, "• Framework base.fixture becomes the single storageState auth source — retire web-pet's fixtures.ts / global-setup.",
         12, WHITE, space_after=3)
    para(tf, "• parent-picker-helpers → reusable Modal / Form / custom-picker components.",
         12, WHITE, space_after=3)
    para(tf, "• Route every spec through one reporting pipeline (Allure / Email / Slack) under CI triggers.",
         12, WHITE, space_after=3)
    para(tf, "• Keep the good patterns as conventions: resolve-ids-by-name + equiv/ parity + JSDoc headers.",
         12, RGBColor(0xC9, 0xD6, 0xE2), space_after=0)
    notes(slide, "This is the slide the architect will probe. Message: we do NOT keep two harnesses. "
                 "The 56 web-pet specs migrate into the framework and reuse its auth, config, and reporting. "
                 "employee.spec.ts (their documented reference, 14/14) is a natural first spec to port. "
                 "We preserve what's valuable — resolve-ids-by-name, the equiv/ parity suite, and the JSDoc "
                 "header convention. Net result: one config, one auth, one reporting pipeline.")

# ================================================================ SLIDE 12 — Roadmap & status
def s_roadmap():
    slide = base_slide("Status & Next Steps")
    top = Inches(1.35)
    cw = Inches(5.95)
    # Framework & Reporting
    rect(slide, MARGIN, top, cw, Inches(4.7), CARD, line=BORDER, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = textbox(slide, MARGIN + Inches(0.3), top + Inches(0.2), cw - Inches(0.6), Inches(4.3))
    para(tf, "Framework & Reporting", 15, NAVY, bold=True, space_after=8, first=True)
    for text, done in [
        ("POM & reusable component architecture", True),
        ("Data-driven testing (JSON / CSV direct)", True),
        ("Fixtures & lifecycle + metrics tracking", True),
        ("API testing with multi-auth strategies", True),
        ("Structured logging & observability", True),
        ("Rich Email & Slack + Allure size-fix", True),
        ("Broaden spec coverage beyond login", False),
    ]:
        p = tf.add_paragraph(); p.space_after = Pt(6)
        r = p.add_run(); r.text = ("✓  " if done else "◷  ") + text
        _set(r, 12.5, GREEN if done else AMBER, bold=done)
    # CI/CD & Migration
    rect(slide, MARGIN + cw + Inches(0.2), top, cw, Inches(4.7), CARD, line=BORDER,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf2 = textbox(slide, MARGIN + cw + Inches(0.5), top + Inches(0.2), cw - Inches(0.6), Inches(4.3))
    para(tf2, "CI/CD & Migration", 15, NAVY, bold=True, space_after=8, first=True)
    for text, done in [
        ("Pipeline (install, browsers, test run)", True),
        ("Triggers: push, cron, manual, dispatch", True),
        ("Reporting & artifacts + Allure trend", True),
        ("Environment config & secrets", True),
        ("ELK dashboard — connect live endpoint", False),
        ("Migrate web-pet e2e (56 specs) onto platform", False),
        ("Preserve equiv/ parity + port picker helpers", False),
    ]:
        p = tf2.add_paragraph(); p.space_after = Pt(6)
        r = p.add_run(); r.text = ("✓  " if done else "◷  ") + text
        _set(r, 12.5, GREEN if done else AMBER, bold=done)
    # legend
    ltf = textbox(slide, MARGIN, Inches(6.3), CONTENT_W, Inches(0.4))
    p = ltf.paragraphs[0]
    r = p.add_run(); r.text = "✓  Done      "; _set(r, 12, GREEN, bold=True)
    r = p.add_run(); r.text = "◷  In progress / next"; _set(r, 12, AMBER, bold=True)
    notes(slide, "Close with an honest status: the framework capabilities, reporting, and the CI/CD pipeline "
                 "are in place; the open items are (1) wiring ELK to a live endpoint, (2) migrating the 56 "
                 "web-pet e2e specs onto the platform while preserving the equiv/ parity suite and porting the "
                 "picker helpers to components, and (3) broadening coverage beyond the login module.")

# ---------------------------------------------------------------- build
s_title()
s_migration()
s_agenda()
s_glance()
s_architecture()
s_structure()
s_cap1()
s_cap2()
s_ai()
s_reporting()
s_cicd_flow()
s_cicd_features()
s_consolidation()
s_roadmap()

out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "Playwright-E2E-Framework-Architecture.pptx")
prs.save(out)
print("Saved:", out)
print("Slides:", len(prs.slides._sldIdLst))
